"""
Document processing module for extracting and chunking text from various file formats.
Supports: .txt, .docx, .pptx, .pdf
"""
import os
from pathlib import Path
from typing import List, Dict, Any
import docx
from pptx import Presentation
from pypdf import PdfReader


class DocumentProcessor:
    """Handles document text extraction and chunking."""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize document processor.

        Args:
            chunk_size: Maximum characters per chunk
            chunk_overlap: Number of overlapping characters between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Process a document file and return extracted text with metadata.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary containing document metadata and chunks

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Extract text based on file type
        extension = path.suffix.lower()
        if extension == '.txt':
            text = self._extract_text_from_txt(path)
        elif extension == '.docx':
            text = self._extract_text_from_docx(path)
        elif extension == '.pptx':
            text = self._extract_text_from_pptx(path)
        elif extension == '.pdf':
            text = self._extract_text_from_pdf(path)
        else:
            raise ValueError(f"Unsupported file format: {extension}")

        # Create chunks
        chunks = self._create_chunks(text)

        return {
            'filename': path.name,
            'file_path': str(path),
            'file_type': extension,
            'full_text': text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'total_chars': len(text)
        }

    def _extract_text_from_txt(self, path: Path) -> str:
        """Extract text from a plain text file."""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()

    def _extract_text_from_docx(self, path: Path) -> str:
        """Extract text from a Word document."""
        doc = docx.Document(path)
        paragraphs = [paragraph.text for paragraph in doc.paragraphs]
        return '\n'.join(paragraphs)

    def _extract_text_from_pptx(self, path: Path) -> str:
        """Extract text from a PowerPoint presentation."""
        prs = Presentation(path)
        text_runs = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide marker
            text_runs.append(f"\n--- Slide {slide_num} ---\n")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)

        return '\n'.join(text_runs)

    def _extract_text_from_pdf(self, path: Path) -> str:
        """Extract text from a PDF document."""
        reader = PdfReader(path)
        text_runs = []

        for page_num, page in enumerate(reader.pages, 1):
            # Add page marker
            text_runs.append(f"\n--- Page {page_num} ---\n")

            # Extract text from the page
            text = page.extract_text()
            if text:
                text_runs.append(text)

        return '\n'.join(text_runs)

    def _create_chunks(self, text: str) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks.

        Args:
            text: Full text to chunk

        Returns:
            List of chunk dictionaries with text and metadata
        """
        if not text or len(text) == 0:
            return []

        chunks = []
        start = 0

        while start < len(text):
            end = start + self.chunk_size

            # If this is not the last chunk, try to break at a sentence or word boundary
            if end < len(text):
                # Look for sentence end (., !, ?) within the last 100 chars
                sentence_end = text.rfind('.', start, end)
                if sentence_end == -1:
                    sentence_end = text.rfind('!', start, end)
                if sentence_end == -1:
                    sentence_end = text.rfind('?', start, end)

                # If found, use sentence boundary
                if sentence_end != -1 and sentence_end > start + (self.chunk_size // 2):
                    end = sentence_end + 1
                else:
                    # Otherwise, try to break at word boundary
                    space_pos = text.rfind(' ', start, end)
                    if space_pos != -1 and space_pos > start:
                        end = space_pos

            chunk_text = text[start:end].strip()

            if chunk_text:  # Only add non-empty chunks
                chunks.append({
                    'text': chunk_text,
                    'start_char': start,
                    'end_char': end,
                    'chunk_index': len(chunks)
                })

            # Move start position with overlap
            start = end - self.chunk_overlap if end < len(text) else end

            # Prevent infinite loop
            if start >= len(text):
                break

        return chunks

    def validate_file(self, file_path: str) -> tuple[bool, str]:
        """
        Validate if a file can be processed.

        Args:
            file_path: Path to the file

        Returns:
            Tuple of (is_valid, error_message)
        """
        path = Path(file_path)

        if not path.exists():
            return False, "File does not exist"

        if not path.is_file():
            return False, "Path is not a file"

        extension = path.suffix.lower()
        supported_formats = ['.txt', '.docx', '.pptx', '.pdf']

        if extension not in supported_formats:
            return False, f"Unsupported format. Supported: {', '.join(supported_formats)}"

        # Check file permissions
        if not os.access(path, os.R_OK):
            return False, "File is not readable (permission denied)"

        return True, ""
